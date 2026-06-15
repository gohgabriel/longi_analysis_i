from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x35, 0x5e)   # headings / title bg
BLUE   = RGBColor(0x2b, 0x6c, 0xb0)   # accent / bullets
TEAL   = RGBColor(0x0d, 0x93, 0x88)   # section dividers
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGREY  = RGBColor(0xf4, 0xf6, 0xf9)   # slide background
DGREY  = RGBColor(0x44, 0x55, 0x66)   # body text
AMBER  = RGBColor(0xd9, 0x7f, 0x06)   # callout accent

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()                    # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=DGREY,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=16, bold=False, color=DGREY,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def slide_bg(slide, color=LGREY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide builders ───────────────────────────────────────────────────────────

def title_slide(title, subtitle=""):
    sl = prs.slides.add_slide(BLANK)
    # Full background
    add_rect(sl, 0, 0, W, H, NAVY)
    # Accent bar
    add_rect(sl, 0, H - Inches(1.2), W, Inches(1.2), TEAL)
    # Title
    add_text_box(sl, title,
                 Inches(1), Inches(2.2), Inches(11.3), Inches(2),
                 font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(sl, subtitle,
                     Inches(1), Inches(4.4), Inches(11.3), Inches(0.8),
                     font_size=22, color=RGBColor(0xb0, 0xc8, 0xe8),
                     align=PP_ALIGN.CENTER)
    return sl


def section_slide(section_title, section_subtitle=""):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LGREY)
    add_rect(sl, 0, 0, Inches(0.25), H, TEAL)   # left accent bar
    add_text_box(sl, section_title,
                 Inches(0.7), Inches(2.8), Inches(11.9), Inches(1.2),
                 font_size=36, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if section_subtitle:
        add_text_box(sl, section_subtitle,
                     Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.8),
                     font_size=20, color=BLUE, align=PP_ALIGN.LEFT)
    return sl


def content_slide(heading, bullets, notes=""):
    """Standard bullet slide. bullets = list of (text, indent_level) tuples."""
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)

    # Header bar
    add_rect(sl, 0, 0, W, Inches(1.1), NAVY)
    add_text_box(sl, heading,
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.8),
                 font_size=26, bold=True, color=WHITE)

    # Accent strip under header
    add_rect(sl, 0, Inches(1.1), W, Inches(0.06), TEAL)

    # Body text box
    txb = slide.shapes.add_textbox if False else \
          sl.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.7))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for (text, level) in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(5 if level == 0 else 2)
        run = p.add_run()
        prefix = "• " if level == 0 else "  – "
        run.text = prefix + text
        run.font.size  = Pt(19 if level == 0 else 16)
        run.font.bold  = (level == 0)
        run.font.color.rgb = NAVY if level == 0 else DGREY

    if notes:
        sl.notes_slide.notes_text_frame.text = notes

    return sl


def two_col_slide(heading, left_title, left_items, right_title, right_items):
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)

    add_rect(sl, 0, 0, W, Inches(1.1), NAVY)
    add_text_box(sl, heading,
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.8),
                 font_size=26, bold=True, color=WHITE)
    add_rect(sl, 0, Inches(1.1), W, Inches(0.06), TEAL)

    col_w = Inches(5.9)
    gap   = Inches(0.5)
    lx    = Inches(0.5)
    rx    = lx + col_w + gap
    cy    = Inches(1.4)
    ch    = Inches(5.6)

    # Left column box
    add_rect(sl, lx, cy, col_w, ch, WHITE)
    add_text_box(sl, left_title, lx + Inches(0.15), cy + Inches(0.1),
                 col_w - Inches(0.3), Inches(0.45),
                 font_size=17, bold=True, color=BLUE)

    txb = sl.shapes.add_textbox(lx + Inches(0.15), cy + Inches(0.65),
                                col_w - Inches(0.3), ch - Inches(0.85))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for txt in left_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run(); run.text = "• " + txt
        run.font.size = Pt(16); run.font.color.rgb = DGREY

    # Right column box
    add_rect(sl, rx, cy, col_w, ch, WHITE)
    add_text_box(sl, right_title, rx + Inches(0.15), cy + Inches(0.1),
                 col_w - Inches(0.3), Inches(0.45),
                 font_size=17, bold=True, color=BLUE)

    txb2 = sl.shapes.add_textbox(rx + Inches(0.15), cy + Inches(0.65),
                                 col_w - Inches(0.3), ch - Inches(0.85))
    txb2.word_wrap = True
    tf2 = txb2.text_frame; tf2.word_wrap = True
    first = True
    for txt in right_items:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run(); run.text = "• " + txt
        run.font.size = Pt(16); run.font.color.rgb = DGREY

    return sl


def apces_row_slide(heading, rows):
    """rows = list of (component, description) — one APCES summary table."""
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)

    add_rect(sl, 0, 0, W, Inches(1.1), NAVY)
    add_text_box(sl, heading,
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.8),
                 font_size=26, bold=True, color=WHITE)
    add_rect(sl, 0, Inches(1.1), W, Inches(0.06), TEAL)

    row_h  = Inches(0.88)
    col1_w = Inches(2.2)
    col2_w = Inches(10.3)
    lx     = Inches(0.35)
    colors = [LGREY, WHITE]

    for i, (comp, desc) in enumerate(rows):
        ry = Inches(1.25) + i * row_h
        bg = colors[i % 2]
        add_rect(sl, lx, ry, col1_w, row_h, NAVY if i == 0 else bg)
        add_rect(sl, lx + col1_w, ry, col2_w, row_h, bg)

        tc = WHITE if i == 0 else BLUE
        add_text_box(sl, comp,
                     lx + Inches(0.1), ry + Inches(0.18),
                     col1_w - Inches(0.2), row_h - Inches(0.2),
                     font_size=15, bold=True, color=tc)
        add_text_box(sl, desc,
                     lx + col1_w + Inches(0.15), ry + Inches(0.12),
                     col2_w - Inches(0.25), row_h - Inches(0.15),
                     font_size=14, color=DGREY)

    return sl


# ════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ════════════════════════════════════════════════════════════════════════════

# ── TITLE ────────────────────────────────────────────────────────────────────
title_slide(
    "Longitudinal Data Analysis",
    "A practical course for programme evaluators"
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 1 — BASICS
# ════════════════════════════════════════════════════════════════════════════
section_slide("Section 1", "Basics of Longitudinal Data")

# 1.1 What is longitudinal data?
content_slide(
    "What is Longitudinal Data?",
    [
        ("Data collected from the same subjects at multiple time points", 0),
        ("Contrast with cross-sectional data — different people, one snapshot", 1),
        ("Key feature: within-person variation across time", 1),
        ("Why it matters for programme evaluation", 0),
        ("A child measured before and after an intervention is not two independent observations", 1),
        ("Within-person correlation is signal, not noise — and must be accounted for", 1),
        ("Core challenge of all methods in this course", 0),
        ("OLS assumes observations are independent — longitudinal data structurally violates this", 1),
        ("Getting standard errors wrong means p-values and CIs cannot be trusted", 1),
    ]
)

# 1.2 The APCES Framework — overview
section_slide("Chapter 1.2", "The APCES Framework")

content_slide(
    "APCES: A Unified Lens for Any Statistical Method",
    [
        ("Every data analysis method can be understood through five components", 0),
        ("A — Architecture: your assumptions about how the world generates data", 1),
        ("P — Parameters: the unknown quantities you want to estimate", 1),
        ("C — Criterion: the scoring rule that defines 'good' estimation", 1),
        ("E — Estimator: the formula that maps data to parameter values", 1),
        ("S — Solver: the computational procedure that actually finds the answer", 1),
        ("Why APCES?", 0),
        ("Change one component and the whole method changes", 1),
        ("Makes it easy to compare OLS, GLS, multilevel models, and Bayes side by side", 1),
    ]
)

content_slide(
    "A — Architecture",
    [
        ("The working theory of how your data was generated", 0),
        ("Written as an equation: e.g.  y = β₀ + β₁·time + ε", 1),
        ("Specifies functional form and how errors behave", 1),
        ("Key: the true parameter θ_true is never directly observable", 0),
        ("We reverse-engineer it from sample data", 1),
        ("Different architectures → different estimation strategies", 1),
        ("Example: assuming ε ~ N(0, σ²) means errors are independent and equally spread", 0),
    ]
)

content_slide(
    "P — Parameters",
    [
        ("The unknowns that, once pinned down, fully describe the model", 0),
        ("θ_true: the real population value (never observed)", 1),
        ("θ̂ (theta-hat): our estimate from the sample", 1),
        ("The hat is a reminder — this is a best guess, not the truth", 1),
        ("Three flavours of parameters to report", 0),
        ("Point estimates — e.g. 'the programme improves scores by 4.2 points'", 1),
        ("Uncertainty estimates — standard errors; how much would our guess vary across samples?", 1),
        ("Covariance structure — how do our estimates relate to each other?", 1),
    ]
)

content_slide(
    "C — Criterion & E — Estimator",
    [
        ("Criterion: the scoring rule defining what 'good' estimation looks like", 0),
        ("OLS criterion: minimise the Residual Sum of Squares (RSS = Σ(y − ŷ)²)", 1),
        ("Also called the 'loss function' in machine learning", 1),
        ("Change the criterion → change what your estimates optimise for", 1),
        ("Estimator: the formula that maps data to a concrete estimate", 0),
        ("OLS estimator: β̂ = (X′X)⁻¹ X′y", 1),
        ("Comes with guarantees: unbiasedness, consistency, efficiency", 1),
        ("Guarantees only hold when architectural assumptions are met (Gauss-Markov)", 1),
    ]
)

content_slide(
    "S — Solver",
    [
        ("The computational procedure that actually finds the answer", 0),
        ("Even with a perfect criterion and estimator, we need something to navigate to the solution", 1),
        ("OLS: analytical solution — one formula, no iteration needed", 0),
        ("More complex models require iterative approaches", 0),
        ("EM algorithm, Newton-Raphson, BOBYQA (for multilevel models)", 1),
        ("MCMC / Variational Inference (for Bayesian models)", 1),
        ("A bad solver can converge to the wrong answer — always check convergence warnings", 1),
    ]
)

apces_row_slide(
    "APCES Applied to OLS — Summary",
    [
        ("Architecture",  "y = β₀ + β₁·time + ε,  ε ~ N(0, σ²). Errors independent, identically distributed."),
        ("Parameters",    "β₀ (baseline), β₁ (treatment effect), σ² (error spread). Estimated as β̂₀, β̂₁, σ̂²."),
        ("Criterion",     "Minimise RSS = Σ(yᵢ − ŷᵢ)²"),
        ("Estimator",     "β̂ = (X′X)⁻¹ X′y — BLUE under Gauss-Markov when four key assumptions hold."),
        ("Solver",        "Analytical: normal equations (X′X)β̂ = X′y. No iteration required."),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 2 — FREQUENTIST METHODS
# ════════════════════════════════════════════════════════════════════════════
section_slide("Section 2", "Frequentist Methods")

# 2.1 Frequentist inference
content_slide(
    "Frequentist Inference — Core Ideas",
    [
        ("The frequentist framework imagines repeating the study many times", 0),
        ("Each repetition gives a slightly different estimate β̂₁", 1),
        ("The distribution of those estimates is the sampling distribution", 1),
        ("p-value: if the null were true, how often would we see a result this extreme?", 0),
        ("Not the probability the null is true — a common misconception", 1),
        ("Confidence interval: the range that would contain θ_true 95% of the time across repetitions", 0),
        ("Not 'there is a 95% chance the true value is in this range'", 1),
        ("Both p-values and CIs depend critically on having correct standard errors", 1),
    ]
)

# 2.2 Naïve OLS
section_slide("Chapter 2.2", "Naïve OLS")

content_slide(
    "Naïve OLS — Three Situations",
    [
        ("Three common designs all reduce to the same equation: y = β₀ + β₁·time + ε", 0),
        ("Situation 1 — Two separate cohorts at different time points", 1),
        ("Different children at t1 vs t2; observations are genuinely independent", 1),
        ("Situation 2 — Two groups at the same time point (intervention vs control)", 1),
        ("Group membership and time are interchangeable", 1),
        ("Situation 3 — Same children measured twice, but treated as independent", 1),
        ("This is the problematic case: OLS estimates β̂₁ correctly, but standard errors are wrong", 1),
    ]
)

two_col_slide(
    "Naïve OLS — When It Works vs When It Breaks",
    "When OLS is valid",
    [
        "Situations 1 & 2: observations are truly independent",
        "Standard errors are correct → CIs and p-values can be trusted",
        "β̂₁ is BLUE (Best Linear Unbiased Estimator) under Gauss-Markov",
        "Comparing groups or cohorts without repeated measures is fine",
    ],
    "When OLS breaks down",
    [
        "Situation 3: same child measured at multiple time points",
        "Within-person correlation violates the independence assumption",
        "β̂₁ is still a reasonable effect estimate, but SEs are too small",
        "p-values too low, CIs too narrow → false confidence",
        "Fix: GLS or multilevel models",
    ]
)

# 2.3 GLS
section_slide("Chapter 2.3", "Generalised Least Squares (GLS)")

content_slide(
    "GLS — The Core Idea",
    [
        ("GLS is the minimal fix to OLS for longitudinal data", 0),
        ("Keeps the same equation: y = β₀ + β₁·time + ε", 1),
        ("Changes exactly one thing: the assumption about how errors behave", 1),
        ("OLS: ε ~ N(0, σ²) — one number, assumes independence", 0),
        ("GLS: ε ~ N(0, Σ) — a matrix that can encode correlation", 0),
        ("Off-diagonal entries of Σ describe the relationship between observations", 1),
        ("Two obs from the same child can be positively correlated", 1),
        ("Two obs from different children stay uncorrelated", 1),
    ]
)

content_slide(
    "GLS — Three Covariance Structures",
    [
        ("We must choose a structure for Σ before estimating", 0),
        ("Compound Symmetry (CS)", 0),
        ("Every pair of obs from the same child is equally correlated", 1),
        ("Simplest assumption; sensible default with only two time points", 1),
        ("AR(1) — First-order Autoregressive", 0),
        ("Observations closer in time are more correlated", 1),
        ("Correlation decays at a fixed rate as the gap grows — more realistic for longer studies", 1),
        ("Unstructured", 0),
        ("Every within-child correlation estimated freely — most flexible, most expensive", 1),
        ("Requires larger samples; compare structures using AIC / BIC", 1),
    ]
)

content_slide(
    "GLS — Estimator & Payoff",
    [
        ("GLS estimator: β̂_GLS = (X′Σ⁻¹X)⁻¹ X′Σ⁻¹y", 0),
        ("Same shape as OLS, but with Σ⁻¹ weighting inserted in two places", 1),
        ("Noisier observations are down-weighted; correlated obs are not double-counted", 1),
        ("Result: BLUE again — but now for correlated data", 0),
        ("Standard errors are now honest → p-values and CIs can be trusted", 1),
        ("Key caveat", 0),
        ("'Correctly specified' is doing a lot of work", 1),
        ("Wrong choice of covariance structure → SEs are still off (less badly than OLS, but still off)", 1),
        ("Try several structures; compare with AIC / BIC", 1),
    ]
)

apces_row_slide(
    "APCES Applied to GLS — Summary",
    [
        ("Architecture",  "y = β₀ + β₁·time + ε,  ε ~ N(0, Σ). Errors can be correlated within child."),
        ("Parameters",    "β₀, β₁, and Σ (covariance structure: compound symmetry, AR(1), or unstructured)."),
        ("Criterion",     "Minimise weighted RSS: Σ(yᵢ − ŷᵢ)′ Σ⁻¹ (yᵢ − ŷᵢ). In practice: maximise REML likelihood."),
        ("Estimator",     "β̂_GLS = (X′Σ⁻¹X)⁻¹ X′Σ⁻¹y — BLUE when Σ correctly specified."),
        ("Solver",        "Iterative EM / Newton-Raphson: estimate β and Σ in alternating steps until convergence."),
    ]
)

# 2.4 Multilevel models
section_slide("Chapter 2.4", "Multilevel Models")

content_slide(
    "Multilevel Models — The Key Intuition",
    [
        ("GLS models the correlation directly; multilevel models explain why it exists", 0),
        ("The orange tree analogy", 0),
        ("Two slices from the same orange taste alike — not mysterious, just the same fruit", 1),
        ("Two measurements from the same child are similar — same underlying child", 1),
        ("Once we label which orange each slice came from, within-orange variation is independent", 1),
        ("Model equation: yᵢⱼ = β₀ + uⱼ + β₁·timeᵢⱼ + εᵢⱼ", 0),
        ("uⱼ = child-specific offset (which 'orange' is this child?)", 1),
        ("εᵢⱼ = residual noise within that child's measurements", 1),
    ]
)

content_slide(
    "Multilevel Models — Variance Decomposition",
    [
        ("Variance is split into two buckets", 0),
        ("Between-children variance τ² — orange-to-orange differences (baseline differences)", 1),
        ("Within-child variance σ² — slice-to-slice differences within the same child", 1),
        ("The ICC (Intraclass Correlation Coefficient) summarises this", 0),
        ("ICC = τ² / (τ² + σ²)", 1),
        ("ICC = 0.6 means 60% of variation is stable differences between children", 1),
        ("High ICC → children are very different from each other; a few obs per child carry little info", 1),
        ("ICC ≈ ρ in GLS under compound symmetry — the two approaches are closely related", 1),
    ]
)

content_slide(
    "Multilevel Models — Estimator & Shrinkage",
    [
        ("Fixed effects β₀, β₁ estimated via weighted regression (like GLS)", 0),
        ("Random effects uⱼ predicted via BLUP (Best Linear Unbiased Prediction)", 0),
        ("BLUP is a shrinkage estimator — a compromise between:", 1),
        ("The child's own data (trust this more when we have many obs per child)", 1),
        ("The population average (pull toward this when we have few obs)", 1),
        ("Shrinkage is a feature, not a bug — guards against overfitting with sparse data", 1),
        ("Key limitation", 0),
        ("Need enough children to estimate τ² reliably (rough guide: ≥ 20)", 1),
        ("With only two time points, compound symmetry is implied — same as GLS", 1),
    ]
)

apces_row_slide(
    "APCES Applied to Multilevel Models — Summary",
    [
        ("Architecture",  "yᵢⱼ = β₀ + uⱼ + β₁·timeᵢⱼ + εᵢⱼ,  uⱼ ~ N(0,τ²),  εᵢⱼ ~ N(0,σ²). Nested observations."),
        ("Parameters",    "β₀, β₁ (fixed effects); τ² (between-child variance); σ² (within-child variance); ICC = τ²/(τ²+σ²)."),
        ("Criterion",     "REML likelihood over the two-level data structure. Use plain ML only when comparing fixed effects."),
        ("Estimator",     "Fixed effects via weighted regression; random effects via BLUP (shrinkage toward the population mean)."),
        ("Solver",        "Iterative: EM + Newton-Raphson (nlme) or BOBYQA + Nelder-Mead (lme4). Rescale time to 0,1,2,… if convergence fails."),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 3 — BAYESIAN METHODS
# ════════════════════════════════════════════════════════════════════════════
section_slide("Section 3", "Bayesian Methods")

# 3.1 Bayesian inference
content_slide(
    "Bayesian Inference — Core Shift",
    [
        ("Frequentist: parameters are fixed unknowns; data are random", 0),
        ("Bayesian: parameters are also treated as random — they have distributions", 0),
        ("Three key quantities", 0),
        ("Prior p(θ): our belief about the parameter before seeing data", 1),
        ("Likelihood p(data | θ): how probable is the data given a parameter value?", 1),
        ("Posterior p(θ | data): updated belief after combining prior and likelihood", 1),
        ("Bayes' theorem: posterior ∝ likelihood × prior", 0),
        ("The posterior is the answer — a full distribution, not a single number", 1),
    ]
)

content_slide(
    "Bayesian Inference — Credible Intervals vs Confidence Intervals",
    [
        ("Confidence interval (frequentist)", 0),
        ("If we repeated the study many times, 95% of CIs would contain θ_true", 1),
        ("Counterintuitive: does NOT mean 95% chance the true value is in this interval", 1),
        ("Credible interval (Bayesian)", 0),
        ("Given the data and our prior, there is a 95% probability θ lies in this range", 1),
        ("Direct probability statement — the natural interpretation people expect", 1),
        ("Prior choice matters", 0),
        ("Uninformative prior → posterior ≈ likelihood → results similar to frequentist", 1),
        ("Informative prior → pulls posterior toward prior beliefs (useful when data are sparse)", 1),
    ]
)

# 3.2 Hierarchical Bayes
section_slide("Chapter 3.2", "Hierarchical Bayes")

content_slide(
    "Hierarchical Bayes — Extending Multilevel Models",
    [
        ("Multilevel model: β₀, β₁ are fixed unknowns; only uⱼ and εᵢⱼ are random", 0),
        ("Hierarchical Bayes: everything is given a distribution — even β₀ and β₁", 0),
        ("Same equation: yᵢⱼ = β₀ + uⱼ + β₁·timeᵢⱼ + εᵢⱼ", 0),
        ("But now priors are placed on all parameters:", 1),
        ("β₀, β₁ ~ Normal (usually centred near zero; wide enough to let data speak)", 1),
        ("τ, σ ~ Half-t (must be positive; weakly informative)", 1),
        ("Key insight: multilevel models already used a prior on uⱼ ~ N(0, τ²)", 0),
        ("HB just extends that logic to the fixed effects too — fully probabilistic top to bottom", 1),
    ]
)

content_slide(
    "Hierarchical Bayes — Criterion & Estimator",
    [
        ("No single number to minimise — goal is to characterise the full posterior", 0),
        ("Posterior ∝ p(data | all params) × p(all params)", 1),
        ("Every parameter gets a full posterior distribution", 1),
        ("Point estimates: posterior mean or median (report when forced to give one number)", 0),
        ("Uncertainty: 95% credible intervals — direct probability statements", 0),
        ("Random effects uⱼ also get full posteriors with explicit shrinkage uncertainty", 0),
        ("Shrinkage toward zero still applies — but now the uncertainty in that shrinkage is visible", 1),
    ]
)

content_slide(
    "Hierarchical Bayes — Solver (MCMC vs VI)",
    [
        ("The posterior is almost never available in closed form with a hierarchical model", 0),
        ("Two main strategies", 0),
        ("MCMC (Markov Chain Monte Carlo)", 1),
        ("Draws samples from the posterior; given enough samples, arbitrarily accurate", 1),
        ("Most common variant: NUTS (No-U-Turn Sampler) — used in brms/Stan", 1),
        ("Expensive: can take minutes to hours for large datasets", 1),
        ("Variational Inference (VI)", 1),
        ("Approximates the posterior with a simpler distribution (e.g. normal)", 1),
        ("Much faster, but introduces approximation bias", 1),
        ("Use MCMC for trustworthy estimates; VI when speed matters and approximation is acceptable", 1),
    ]
)

content_slide(
    "Hierarchical Bayes — Practical Considerations",
    [
        ("Computational cost", 0),
        ("MCMC can be slow; check convergence diagnostics (R̂ statistic, trace plots)", 1),
        ("A sampler that has not converged produces garbage silently", 1),
        ("Prior sensitivity", 0),
        ("All priors influence the posterior; effect is larger with small samples", 1),
        ("Best practice: run a sensitivity analysis with different priors", 1),
        ("Same small-sample limitations as multilevel models", 0),
        ("With very few children, posterior over τ will be wide and prior-dominated", 1),
        ("More higher-level units (children, schools) always help", 1),
    ]
)

apces_row_slide(
    "APCES Applied to Hierarchical Bayes — Summary",
    [
        ("Architecture",  "yᵢⱼ = β₀ + uⱼ + β₁·timeᵢⱼ + εᵢⱼ, with priors on all parameters. Everything is random."),
        ("Parameters",    "β₀, β₁, τ, σ, and all uⱼ — each is a posterior distribution, not a point estimate."),
        ("Criterion",     "Characterise the full joint posterior: p(params | data) ∝ p(data | params) × p(params)."),
        ("Estimator",     "Posterior mean/median as point estimate; 95% credible intervals for uncertainty."),
        ("Solver",        "MCMC (NUTS) for exact characterisation; Variational Inference for fast approximation. Both in brms (R)."),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# CLOSING — METHODS COMPARISON
# ════════════════════════════════════════════════════════════════════════════
section_slide("Putting It All Together", "Comparing the Four Methods")

content_slide(
    "Method Comparison at a Glance",
    [
        ("Naïve OLS", 0),
        ("Use when: different people at each time point (no repeated measures)", 1),
        ("Breaks when: same people measured multiple times — SEs are wrong", 1),
        ("GLS", 0),
        ("Use when: longitudinal data; you want to specify the correlation structure explicitly", 1),
        ("Key decision: choosing the right Σ structure (CS, AR(1), unstructured)", 1),
        ("Multilevel Models", 0),
        ("Use when: longitudinal data; want to explain the source of correlation", 1),
        ("Advantage: random effects, ICC, shrinkage — more interpretable variance decomposition", 1),
        ("Hierarchical Bayes", 0),
        ("Use when: you have prior knowledge to incorporate, or you want full uncertainty quantification", 1),
        ("Trade-off: computational cost + need to justify prior choices", 1),
    ]
)

content_slide(
    "The Progressive Fix",
    [
        ("All methods share the same core equation: y = β₀ + β₁·time + ε", 0),
        ("Each method is a more principled way of handling the same core problem:", 1),
        ("Within-person correlation violates OLS's independence assumption", 1),
        ("Step 1 → GLS: model the correlation structure directly", 0),
        ("Step 2 → Multilevel: explain why the correlation exists (random effects)", 0),
        ("Step 3 → Hierarchical Bayes: treat everything as uncertain — full probabilistic reasoning", 0),
        ("The APCES framework lets you see exactly what changes at each step", 0),
        ("Solver and architecture are the main things that shift — the estimator logic is preserved", 1),
    ]
)

content_slide(
    "Key Takeaways",
    [
        ("Longitudinal data violates OLS's independence assumption", 0),
        ("The point estimate β̂₁ may still be reasonable, but standard errors are not", 1),
        ("GLS fixes the SEs by modelling correlation explicitly — but you must pick a structure", 0),
        ("Multilevel models do the same by modelling the source of correlation — more interpretable", 0),
        ("Hierarchical Bayes extends MLM by treating all parameters as distributions", 0),
        ("Choose your method based on your data structure, sample size, and inferential goals", 0),
        ("The APCES framework is your map — any time you encounter a new method, ask:", 1),
        ("What does the architecture assume? What does the criterion optimise? How does the solver work?", 1),
    ]
)

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\calyc\PycharmProjects\misc_stats\longitudinal_course.pptx"
prs.save(out)
print(f"Saved: {out}")
